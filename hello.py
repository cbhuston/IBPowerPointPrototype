from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image

#Current to dos in test file 10/1: 
#Figure out image math/formula for centering 
#Figure out exact math/formula for placing strings on varible heights 

prs = Presentation()
#slide 3
title_only_slide_layout = prs.slide_layouts[6]
slide3 = prs.slides.add_slide(title_only_slide_layout)
shapes = slide3.shapes

#all values all currrently hardcoded not loaded 
height = Inches(1.5) 
width = Inches(1)
left = Inches(.5)
top = Inches(1) 

def x(height, width, left, top):
    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes
    title_only_slide_layout = prs.slide_layouts[6]
    slide3 = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide3.shapes
    for y in range(1, 4):
        left = Inches(.5)
        for x in range(1,8):
            #init shapes
            shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            #set box shadow
            shadow = shape.shadow
            shadow.inherit = False
            print(shadow.inherit)
            #set fill and fill color
            fill = shape.fill
            fill.solid()
            fill.fore_color.rgb = RGBColor(255, 255, 255)
            #set outline
            line = shape.line
            line.color.rgb = RGBColor(0,0,0)
            line.width = Pt(2.5)
            
            #setting height of shape
            txBoxHeight = Inches(.25)
            txBox1 = slide3.shapes.add_textbox(left, top-Inches(.33), width, txBoxHeight)
            tf = txBox1.text_frame
            p = tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            #todo: get from csv
            p.text = "June 2018"
            p.font.size = Pt(10)

            picTop = top+Inches(.4)#varible place image further down without effecting top

            im = Image.open('key_social_logo.png')
            picwidth, picheight = im.size #im.size should bring in a width and height
            # image.left = (tombstone_width - image.width) / 2 **Implemented below**
            picleft = ((width - (picheight/picwidth)*Inches(.4))/8) + left
            pic = slide3.shapes.add_picture('key_social_logo.png', picleft, picTop, height=Inches(.4))
            txBoxTop = top + Inches(1.5)# varible to move txBox without effecting top 
            
            txBox3 = slide3.shapes.add_textbox(left, picTop+Inches(.1), width, Inches(.25))
            tf = txBox3.text_frame
            p = tf.add_paragraph()
            p.text = "200,000,000"
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER

            txBox3 = slide3.shapes.add_textbox(left, picTop+Inches(.23), width, Inches(.25))
            tf = txBox3.text_frame
            p = tf.add_paragraph()
            p.text = "IPO"
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER

            txBox2 = slide3.shapes.add_textbox(left, top+Inches(1), width, Inches(.25))
            tf = txBox2.text_frame
            p = tf.add_paragraph()
            p.text = "Bookrunner"
            p.font.size = Pt(10)
            p.alignment = PP_ALIGN.CENTER
            left = left + width + Inches(.285)
        top = top + height + Inches(.33)

    prs.save('test.pptx')
x(height, width, left, top)
