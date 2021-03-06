import os
from flask import Flask, flash, request, redirect, url_for
from werkzeug.utils import secure_filename
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches

app = Flask(__name__)

@app.route('/tombstone/api/v1.0/upload', methods=['POST'])
def upload_data():
    # check if the post request has the file part
  if request.method == 'POST':
    f = request.files['file']
    f.save(secure_filename(f.filename))
    print(f.filename)

    prs = Presentation()
    title_only_slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide.shapes

    shapes.title.text = 'Adding an AutoShape'

    left = Inches(0.93)  # 0.93" centers this overall set of shapes
    top = Inches(3.0)
    width = Inches(1.75)
    height = Inches(1.0)

    shape = shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, width, height)
    shape.text = 'Step 1'

    left = left + width - Inches(0.4)
    width = Inches(2.0)  # chevrons need more width for visual balance

    for n in range(2, 6):
        shape = shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, width, height)
        shape.text = 'Step %d' % n
        left = left + width - Inches(0.4)

    left = top = width = height = Inches(1.0)
    slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )

    title_only_slide_layout = prs.slide_layouts[2]
    slide2 = prs.slides.add_slide(title_only_slide_layout)
    shapes = slide2.shapes

    height = Inches(1.5) 
    width = Inches(1)
    left = Inches(.5)
    top = Inches(1) 

    for y in range(1, 4):
        left = Inches(.5)
        for x in range(1,8):
            shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
            shape.text = 'Deal'
            left = left + width + Inches(.285)
        top = top + height + Inches(.33)

    prs.save('test.pptx')
    return 'file uploaded successfully'

if __name__ == '__main__':
    app.run(debug=True)